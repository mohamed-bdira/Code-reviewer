"""Generate graduation defense PowerPoint from PFE report structure."""

from pptx import Presentation
from pptx.util import Inches, Pt

OUTPUT = r"c:\Users\conta\OneDrive\Documents\Code-reviewer\docs\PFE-Graduation-Presentation-Mohamed-Bdira.pptx"

# Each entry: (title, bullets, speaker_notes)
SLIDES = [
    (
        "Automated Code Governance\nand Pull Request Review System",
        [
            "Presented in order to obtain the",
            "National Bachelor's Degree in Computer Science",
            "",
            "Presented By: Mohamed Bdira",
            "Academic Supervisor: Ms. Noura Aboudi",
            "Professional Supervisor: Mr. Mehdi Mili",
        ],
        "Good morning everyone, and thank you for being here today. My name is Mohamed Bdira, and I will present my graduation project entitled Automated Code Governance and Pull Request Review System. This work was carried out at iMaxeam under the academic supervision of Mrs. Noura Aboudi and the professional supervision of Mr. Mehdi Mili, as part of the Software Engineering and Information Systems program at Horizon School of Digital Technologies.",
    ),
    (
        "Table of Contents",
        [
            "01  Host Company",
            "02  Context & Problem",
            "03  Existing Solutions",
            "04  Proposed Solution",
            "05  Methodology",
            "06  Requirements",
            "07  Architecture & Design",
            "08  Application Demo",
            "09  Conclusion & Perspectives",
        ],
        "This presentation follows the structure of my report. We will begin with the host company, then the context and problem statement, a review of existing solutions, and our proposed platform. After that, I will cover the Agile Scrum methodology, functional and non-functional requirements, system architecture, a live demo walkthrough, and finally conclusions with future perspectives.",
    ),
    (
        "Host Company",
        ["01"],
        "Let us start with the professional context of this project.",
    ),
    (
        "iMaxeam",
        [
            "IBM Gold Business Partner — Enterprise Asset Management",
            "Software publishing, consulting & IT services",
            "Global expertise in IBM Maximo integrations",
            "Sousse, Tunisia — imaxeam.com",
        ],
        "iMaxeam is a prominent IT services and software publishing firm, recognized as an official IBM Gold Business Partner. The company specializes in Enterprise Asset Management, with deep expertise in IBM Maximo integrations, maintenance, and technical architecture optimization. iMaxeam serves asset-intensive industries such as manufacturing, utilities, transportation, and oil and gas. This internship gave me exposure to professional software engineering practices in a real enterprise environment.",
    ),
    (
        "Context & Problem",
        ["02"],
        "",
    ),
    (
        "— Martin Fowler",
        [
            '"Any fool can write code',
            "that a computer can understand.",
            'Good programmers write code',
            'that humans can understand."',
        ],
        "Code review is fundamentally about human understanding. Even with modern CI/CD pipelines automating builds, tests, and linting, the review phase itself remains largely manual. That gap is what this project addresses.",
    ),
    (
        "Key Challenges in Code Review",
        [
            "Engineering Resource Drain",
            "Cognitive Fatigue & Human Error",
            "Lack of Standardization",
        ],
        "Three critical challenges motivated this work. First, senior engineers spend substantial time on repetitive review tasks — style enforcement, trivial syntax fixes, and pattern checking — reducing time for high-value work like system design. Second, pull requests often contain dozens of files and hundreds of lines. Under deadline pressure, reviewer fatigue increases the likelihood that bugs, security vulnerabilities, and technical debt go unnoticed. Third, human reviews reflect individual preferences, leading to inconsistent enforcement of best practices and architectural drift across teams.",
    ),
    (
        "Existing Solutions",
        ["03"],
        "",
    ),
    (
        "Existing Solutions",
        [
            "Solution          Pros                    Cons",
            "SonarQube         Deterministic, fast     No semantic awareness",
            "                  Rule-based              High false positives",
            "",
            "CodeRabbit        Context-aware AI        Closed-source SaaS",
            "                  Conversational PR       No granular AST control",
            "                  reviews                 Data privacy concerns",
        ],
        "We studied two dominant archetypes. SonarQube is the industry standard for static analysis. It is highly optimized and deterministic, excellent at syntax compliance and complexity thresholds. However, it lacks semantic awareness — it reads code as structural tokens, not logical intent — and produces high volumes of false positives. CodeRabbit represents the modern LLM-driven approach. It understands context and intent, summarizes complex changes, and engages developers in PR conversations. But it operates as a closed-source black-box SaaS with per-seat costs and limited ability to combine custom AST rules with AI prompts on a per-repository basis.",
    ),
    (
        "The Gap",
        [
            "SonarQube → Rigid structure, no context",
            "CodeRabbit → Flexible context, no control",
            "",
            "Our platform → Hybrid AST + LLM + Dashboard",
        ],
        "The evaluation reveals a distinct architectural gap. SonarQube provides rigid, deterministic structure without context. CodeRabbit provides flexible context without granular, deterministic control. Our platform fills this gap with a hybrid architecture combining ast-grep structural analysis, Google Gemini semantic reasoning, and a React dashboard for per-repository policy configuration.",
    ),
    (
        "Proposed Solution",
        ["04"],
        "",
    ),
    (
        "AI-Assisted PR Review Platform",
        [
            "Event-driven GitHub webhook integration",
            "Hybrid AST + LLM code analysis",
            "Centralized admin dashboard",
            "Persistent findings & real-time updates",
        ],
        "Our platform is a full-stack, event-driven middleware layer that bridges GitHub with advanced Large Language Models. When a pull request is opened or updated, the system receives a webhook, performs hybrid analysis combining deterministic AST parsing with generative semantic reasoning, persists structured findings in MongoDB, and optionally publishes inline review comments back to GitHub. Operators manage installations, repository policies, and findings through a centralized React dashboard.",
    ),
    (
        "Key Strengths",
        [
            "01  Hybrid Analysis — AST + LLM",
            "02  Per-Repo Configuration",
            "03  GitHub-Native Integration",
            "04  Real-Time Dashboard (SSE)",
        ],
        "Four pillars define our solution. First, hybrid analysis combines ast-grep for deterministic structural checks with Google Gemini for semantic reasoning. Second, per-repository configuration lets team leads define focus areas, enforcement levels, custom rules, and merge score thresholds via RepoConfig. Third, deep GitHub integration through a registered GitHub App with OAuth authentication, signed webhooks, and installation-scoped tokens. Fourth, a real-time dashboard using Server-Sent Events so operators see findings and status updates without manual page reloads.",
    ),
    (
        "Platform in Numbers",
        [
            "5  MongoDB collections",
            "7  Dashboard panels",
            "4  Agile Scrum sprints",
            "1.48s  Mean pipeline latency",
        ],
        "Some key metrics from our implementation and evaluation. The platform uses five MongoDB collections for users, installations, repository configs, API keys, and review findings. The dashboard exposes seven operational panels. Development followed four Agile Scrum sprints from preparation through deployment. In production evaluation, the end-to-end review pipeline achieved a mean latency of 1.48 seconds, well under our 2.5-second threshold.",
    ),
    (
        "Methodology",
        ["05"],
        "",
    ),
    (
        "Agile Scrum",
        [
            "Framework: Agile Scrum",
            "Why Scrum: Evolving AI + GitHub components",
            "Process: Sprint planning → Build → Review",
            "Benefit: Incremental delivery & early refinement",
        ],
        "We used Agile Scrum because the platform combines several evolving components — GitHub webhook integration, AI-assisted analysis, database persistence, and a React dashboard. Instead of a rigid sequential process, we organized work into a product backlog and delivered incrementally across four sprints: Sprint 0 for preparation and requirements, Sprint 1 for backend architecture, Sprint 2 for the AI engine, Sprint 3 for the frontend, and Sprint 4 for deployment and validation. Sprint reviews let us refine LLM response formatting, token limitations, and diff processing early.",
    ),
    (
        "Requirements",
        ["06"],
        "",
    ),
    (
        "Identification of Actors",
        [
            "Developer — Configures repos, opens PRs",
            "GitHub — Webhooks, OAuth, REST API",
            "Google Gemini — Semantic code analysis",
        ],
        "Three primary actors interact with the system. The Developer authenticates via GitHub OAuth, links GitHub App installations, configures repository review policies, and opens or synchronizes pull requests. GitHub acts as an external system delivering webhooks, handling OAuth, and serving REST APIs for diffs and review posting. Google Gemini serves as the AI engine, receiving curated diffs and returning structured review findings.",
    ),
    (
        "Functional Requirements",
        [
            "GitHub OAuth authentication",
            "Installation & repo configuration CRUD",
            "Webhook-driven automated PR review",
            "Findings persistence & dashboard API",
            "SSE live updates & API key management",
        ],
        "Key functional requirements include: GitHub OAuth sign-in with JWT sessions; linking GitHub App installations to operator accounts; CRUD operations on per-repository RepoConfig policies; automated review triggered on pull_request opened and synchronize events; persisting structured findings with deduplication; exposing findings through filterable dashboard APIs; Server-Sent Events for real-time UI updates; and API key lifecycle management for dashboard access and optional review gating.",
    ),
    (
        "Non-Functional Requirements",
        [
            "Security — HMAC, JWT, bcrypt, CORS",
            "Performance — Diff segmentation, retries",
            "Reliability — Early webhook ack, dedupe",
            "Usability — Responsive dashboard, SSE",
        ],
        "Non-functional requirements span four areas. Security includes HMAC-SHA256 webhook validation, JWT sessions, bcrypt-hashed API keys, and single-origin CORS. Performance covers adaptive diff segmentation for large PRs, retry logic for Gemini rate limits, and conservative GitHub posting to avoid secondary rate limits. Reliability means immediate HTTP 200 webhook acknowledgment, finding deduplication via dedupeKey, and MongoDB index enforcement. Usability focuses on a responsive Tailwind CSS dashboard with real-time SSE updates and clear operator workflows.",
    ),
    (
        "Architecture & Design",
        ["07"],
        "",
    ),
    (
        "General System Architecture",
        ["[Insert Figure 3.11 — High-level architecture]"],
        "The architecture follows an event-driven model. GitHub sends pull request webhook events to the Express backend on Railway. The backend validates signatures, resolves tenant context via Installation and RepoConfig, fetches and filters diffs, invokes the Gemini SDK for analysis, persists findings in MongoDB Atlas, and optionally posts review feedback to GitHub. The React SPA on Vercel consumes REST APIs and SSE event streams for monitoring and configuration. This decoupled split keeps the UI fast on a CDN while the API stays always-on for webhooks and background processing.",
    ),
    (
        "Class Diagram",
        ["[Insert Figure 3.2 — Class diagram]"],
        "The data model centers on five MongoDB entities. User stores operator identity from GitHub OAuth. Installation links a user to a GitHub App installation ID. RepoConfig holds per-repository review policy — focus areas, enforcement level, custom rules, merge minimum score, and ast-grep flag. ApiKey stores bcrypt-hashed credentials for dashboard access. PrReviewFinding captures AI-detected issues with category, file path, line range, and dedupe key. The ReviewPullRequest service orchestrates diff fetching, segmentation, Gemini invocation, enforcer parsing, and GitHub posting.",
    ),
    (
        "Webhook Review Sequence",
        ["[Insert Figure 4.4 — Webhook to AI sequence]"],
        "This sequence diagram shows the core business flow. GitHub delivers a pull_request webhook. Express verifies the HMAC signature, loads Installation and RepoConfig, and immediately returns HTTP 200. Then reviewPullRequest fetches the diff via Octokit, filters and segments it, runs ast-grep if enabled, calls Gemini through the official SDK, parses the structured JSON response via the enforcer, upserts findings into MongoDB, posts a GitHub review if appropriate, and publishes an SSE event so the dashboard refreshes automatically.",
    ),
    (
        "Technology Stack",
        [
            "Frontend — React 19, Vite, Tailwind CSS",
            "Backend — Node.js 22, Express 5, TypeScript",
            "Database — MongoDB Atlas, Mongoose",
            "AI — Google Gemini SDK",
            "DevOps — Railway, Vercel, Docker, Git",
        ],
        "The technology stack is organized by operational layer. The frontend is a React 19 single-page application built with Vite and styled with Tailwind CSS, deployed on Vercel. The backend runs Node.js 22 with Express 5 and TypeScript, containerized and hosted on Railway. MongoDB Atlas provides document-oriented persistence through Mongoose ODM. Google Gemini SDK powers semantic analysis with structured JSON output. GitHub integration uses Octokit for App authentication and API calls. Git manages versioning, and CI/CD automates deployment to both hosting platforms.",
    ),
    (
        "Hybrid Analysis Pipeline",
        [
            "1. Webhook ingress & validation",
            "2. Diff fetch, filter & segment",
            "3. ast-grep structural check",
            "4. Gemini semantic analysis",
            "5. Enforcer parse & persist",
            "6. GitHub review publication",
        ],
        "The review pipeline has six phases. First, GitHub PR events trigger webhook ingress with cryptographic validation. Second, the orchestrator fetches the unified diff, strips lockfiles and binaries, and segments large patches. Third, ast-grep runs optional structural pattern matching. Fourth, Gemini analyzes semantic intent with repository policy injected into the prompt. Fifth, the enforcer validates JSON output, computes merge readiness scores, and upserts findings. Sixth, Octokit publishes inline comments and review summaries back to the PR thread.",
    ),
    (
        "Security Measures",
        [
            "HMAC-SHA256 webhook verification",
            "Short-lived installation tokens (60 min)",
            "bcrypt API key hashing (cost 12)",
            "RSA private key in volatile memory only",
        ],
        "Security follows a defense-in-depth approach. Every webhook payload is verified with HMAC-SHA256 using a shared secret and timing-safe comparison. GitHub App RSA private keys are injected via environment variables and kept in volatile memory only. Installation access tokens are short-lived with a 60-minute TTL, requested on demand. Operator API keys are bcrypt-hashed with cost factor 12 and never stored in plaintext. CORS restricts API access to the configured frontend origin.",
    ),
    (
        "Application Demo",
        ["08"],
        "",
    ),
    (
        "Application Demo — Login & GitHub Setup",
        ["[Insert screenshot — GitHub OAuth & App install]"],
        "For the demo, the operator starts at the login page and signs in with GitHub OAuth. After authentication, they navigate to the GitHub and CI panel to link a GitHub App installation. The platform validates ownership via Octokit and stores the installation binding. From there, the operator selects repositories and creates RepoConfig entries with focus areas, enforcement level, custom rules, and optional ast-grep enablement.",
    ),
    (
        "Application Demo — Review Flow",
        ["[Insert screenshot — Findings panel & PR review]"],
        "When a developer opens or updates a pull request on a configured repository, the webhook triggers automated review. Within about 1.5 seconds, findings appear in the dashboard Findings panel via SSE without manual reload. The operator can filter by repository, category, or severity. Simultaneously, inline review comments may appear on the GitHub PR thread. The Overview panel shows aggregated statistics — total findings, category distribution, and service health.",
    ),
    (
        "Deployment Architecture",
        [
            "Frontend → Vercel (CDN, static SPA)",
            "Backend → Railway (Docker, always-on)",
            "Database → MongoDB Atlas (cloud)",
            "CI/CD → Auto-deploy on main push",
        ],
        "Production uses a split deployment. The React SPA builds with Vite and deploys to Vercel for global CDN delivery and HTTPS. The Express API runs in a Docker container on Railway, staying always-on for webhooks, SSE connections, scheduled scans, and Gemini API calls. MongoDB Atlas hosts the database. Environment variables manage secrets — JWT secret, GitHub App key, webhook secret, Gemini API key, and MongoDB URI. CI/CD triggers automatic rebuilds on push to main for both platforms.",
    ),
    (
        "Conclusion & Perspectives",
        ["09"],
        "",
    ),
    (
        "Conclusion",
        [
            "Problem — Manual review bottlenecks CI/CD",
            "Solution — Hybrid AST + LLM platform",
            "Delivered — Production-ready full-stack app",
            "Impact — Faster feedback, consistent standards",
        ],
        "In conclusion, this project addressed the bottleneck of manual code review in modern CI/CD pipelines. We designed and implemented a hybrid platform combining ast-grep structural analysis with Google Gemini semantic reasoning, integrated natively into GitHub workflows. The result is a production-ready full-stack application with a React dashboard, Express backend, MongoDB persistence, and split cloud deployment. The platform delivers near-instant feedback on pull requests while keeping human reviewers in control of merge decisions.",
    ),
    (
        "Perspectives",
        [
            "Short-term — Message queue (Redis/RabbitMQ)",
            "Medium-term — Vector filtering, extended AST",
            "Long-term — Multi-model support, on-premise",
        ],
        "Future work spans three horizons. In the short term, integrating a distributed message queue like Redis or RabbitMQ will isolate webhook ingestion from analysis execution during high-concurrency bursts. In the medium term, local vector filtering can reduce unnecessary Gemini API calls for cosmetic changes, and extended AST compilation will further reduce hallucinations. In the long term, configurable multi-model support per repository — Gemini, GPT, or open-source models — and a fully containerized on-premise deployment blueprint will serve organizations with strict data privacy requirements.",
    ),
    (
        "Thank You!",
        ["Do you have any questions?"],
        "Thank you for your attention. I am happy to answer any questions about the architecture, the AI integration, the GitHub workflow, or the evaluation results. Thank you to my supervisors Mrs. Noura Aboudi and Mr. Mehdi Mili, the iMaxeam team, and the jury members for their time and feedback.",
    ),
]


def add_slide(prs, title, bullets, notes):
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(20 if len(bullets) <= 4 else 16)

    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes


def add_section_slide(prs, title, section_num, notes=""):
    layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = f"{title}\n{section_num}"
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Automated Code Governance\nand Pull Request Review System"
    slide.placeholders[1].text = (
        "Presented in order to obtain the National Bachelor's Degree in Computer Science\n\n"
        "Presented By: Mohamed Bdira\n"
        "Academic Supervisor: Ms. Noura Aboudi  |  Professional Supervisor: Mr. Mehdi Mili"
    )
    slide.notes_slide.notes_text_frame.text = SLIDES[0][2]

    content_slides = SLIDES[1:]
    section_nums = {"01", "02", "03", "04", "05", "06", "07", "08", "09"}

    for title, bullets, notes in content_slides:
        if len(bullets) == 1 and bullets[0] in section_nums:
            s = prs.slides.add_slide(prs.slide_layouts[5])
            s.shapes.title.text = f"{title}\n{bullets[0]}"
            if notes:
                s.notes_slide.notes_text_frame.text = notes
        else:
            add_slide(prs, title, bullets, notes)

    prs.save(OUTPUT)
    print(f"Created {len(prs.slides)} slides -> {OUTPUT}")


if __name__ == "__main__":
    main()
